import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup

from haystack.components.fetchers.link_content import LinkContentFetcher
from haystack.components.converters import HTMLToDocument
from haystack.components.preprocessors import DocumentSplitter
from haystack.components.rankers import TransformersSimilarityRanker
from haystack.components.builders.prompt_builder import PromptBuilder
from haystack import Pipeline
from haystack.components.generators import HuggingFaceLocalGenerator

class QAChatbot:
    def __init__(self):
        self.pipeline = self._create_pipeline()

    def _create_pipeline(self):
        # web_retriever = WebRetriever()
        fetcher = LinkContentFetcher()
        converter = HTMLToDocument()
        document_splitter = DocumentSplitter(split_by="word", split_length=50)
        similarity_ranker = TransformersSimilarityRanker(top_k=3)
        
        generator = HuggingFaceLocalGenerator(
            model="google/flan-t5-large",
            task="text2text-generation",
            # device=ComponentDevice.from_str("cpu"),  # Use GPU if available, otherwise use "cpu"
            generation_kwargs={
                "max_new_tokens": 500, 
                "temperature": 0.7
            }
        )
        generator.warm_up()

        prompt_template = """
        According to these documents:

        {% for doc in documents %}
          {{ doc.content }}
        {% endfor %}

        Answer the given question: {{question}}
        Answer:
        """
        prompt_builder = PromptBuilder(template=prompt_template)

        pipeline = Pipeline()
        pipeline.add_component("fetcher", fetcher)
        pipeline.add_component("converter", converter)
        pipeline.add_component("splitter", document_splitter)
        pipeline.add_component("ranker", similarity_ranker)
        pipeline.add_component("prompt_builder", prompt_builder)
        pipeline.add_component("generator", generator)

        pipeline.connect("fetcher.streams", "converter.sources")
        pipeline.connect("converter.documents", "splitter.documents")
        pipeline.connect("splitter.documents", "ranker.documents")
        pipeline.connect("ranker.documents", "prompt_builder.documents")
        pipeline.connect("prompt_builder.prompt", "generator")

        return pipeline

    def ask_question(self, question, urls):
        result = self.pipeline.run({
            "prompt_builder": {"question": question},
            "ranker": {"query": question},
            "fetcher": {"urls": urls}
        })

        for answer in result["generator"]["replies"]:
            print(answer)
        
        return result["generator"]["replies"]

    def ingest_pdf(self, file_path):
        docs = []
        pdf_reader = PdfReader(file_path)
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text.strip():
                docs.append({"content": text, "meta": {"filename": os.path.basename(file_path)}})
        return docs

    def ingest_docx(self, file_path):
        docs = []
        doc = Document(file_path)
        for para in doc.paragraphs:
            if para.text.strip():
                docs.append({"content": para.text, "meta": {"filename": os.path.basename(file_path)}})
        return docs

    def ingest_pptx(self, file_path):
        docs = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    slide_text.append(shape.text)
            if slide_text:
                docs.append({"content": "\n".join(slide_text), "meta": {"filename": os.path.basename(file_path)}})
        return docs

    def ingest_html(self, file_path):
        docs = []
        with open(file_path, 'r', encoding='utf-8') as file:
            soup = BeautifulSoup(file, 'lxml')
            text = soup.get_text(separator='\n', strip=True)
            docs.append({"content": text, "meta": {"filename": os.path.basename(file_path)}})
        return docs

    def ingest_documents(self, directory):
        docs = []
        for filename in os.listdir(directory):
            file_path = os.path.join(directory, filename)
            if filename.endswith('.pdf'):
                docs.extend(self.ingest_pdf(file_path))
            elif filename.endswith('.docx'):
                docs.extend(self.ingest_docx(file_path))
            elif filename.endswith('.pptx'):
                docs.extend(self.ingest_pptx(file_path))
            elif filename.endswith('.html'):
                docs.extend(self.ingest_html(file_path))
        
        if docs:
            self.add_documents(docs)
            print(f"Ingested {len(docs)} documents from {len(os.listdir(directory))} files.")
        else:
            print("No valid documents found in the specified directory.")

def main():
    chatbot = QAChatbot()
    
    print("Welcome to the Q&A Chatbot! Type 'exit' to quit.")
    while True:
        user_input = input("You: ")
        if user_input.lower() == 'exit':
            break
        urls = input("Enter URLs to search (comma-separated, or press enter to skip): ").split(',')
        urls = [url.strip() for url in urls if url.strip()]
        if not urls:
            urls = ["https://haystack.deepset.ai/blog/introducing-haystack-2-beta-and-advent"]  
        
        answers = chatbot.ask_question(user_input, urls)
        print("Chatbot:")
        for answer in answers:
            print(answer)

if __name__ == "__main__":
    main()